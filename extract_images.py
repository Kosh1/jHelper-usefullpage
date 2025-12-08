# extract_images.py
from pptx import Presentation
import os
import json

def extract_images_from_pptx(pptx_path, output_dir='extracted_images'):
    """Извлекает изображения из PPTX файла"""
    prs = Presentation(pptx_path)
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    image_data = []
    image_count = 0
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_images = []
        
        for shape in slide.shapes:
            # Проверяем, есть ли изображение
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                
                # Определяем расширение
                ext = image.ext
                if not ext:
                    ext = 'png'  # По умолчанию PNG
                
                # Сохраняем изображение
                image_filename = f"slide_{slide_num}_img_{image_count}.{ext}"
                image_path = os.path.join(output_dir, image_filename)
                
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                slide_images.append({
                    "filename": image_filename,
                    "path": image_path,
                    "slide_number": slide_num
                })
                
                image_count += 1
                print(f"Извлечено изображение: {image_filename} (слайд {slide_num})")
        
        if slide_images:
            image_data.append({
                "slide_number": slide_num,
                "images": slide_images
            })
    
    # Сохраняем метаданные
    with open(os.path.join(output_dir, 'images_metadata.json'), 'w', encoding='utf-8') as f:
        json.dump(image_data, f, ensure_ascii=False, indent=2)
    
    print(f"\nВсего извлечено изображений: {image_count}")
    return image_data

if __name__ == "__main__":
    pptx_file = "Нейросети для адвокатов V1.pptx"
    
    try:
        extract_images_from_pptx(pptx_file)
    except Exception as e:
        print(f"Ошибка: {e}")
        import traceback
        traceback.print_exc()

