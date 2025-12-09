# extract_pptx.py
from pptx import Presentation
import json

def extract_pptx_text(pptx_path):
    """Извлекает весь текст из PPTX файла"""
    prs = Presentation(pptx_path)
    
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        # Извлекаем текст из всех фигур на слайде
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        if slide_text:
            slides_data.append({
                "slide_number": slide_num,
                "text": "\n".join(slide_text),
                "text_parts": slide_text
            })
    
    return slides_data

# Использование
if __name__ == "__main__":
    pptx_file = "Нейросети для адвокатов V1.5.pptx"
    
    try:
        data = extract_pptx_text(pptx_file)
        
        # Сохраняем в JSON
        with open("presentation_data.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        
        # Также выводим в консоль
        print(f"Извлечено {len(data)} слайдов:\n")
        for slide in data:
            print(f"--- Слайд {slide['slide_number']} ---")
            print(slide['text'])
            print()
            
    except ImportError:
        print("Установите библиотеку: pip install python-pptx")
    except Exception as e:
        print(f"Ошибка: {e}")


